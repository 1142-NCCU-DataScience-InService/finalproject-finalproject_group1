"""
產生「猴子都看得懂」的期末報告簡報 PPTX
執行：python generate_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = "車禍精神慰撫金AI預測_期末報告簡報.pptx"

# 配色
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x00, 0x68, 0xC9)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Microsoft JhengHei"
    return tb


def add_box(slide, left, top, width, height, fill=LIGHT, line=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line; box.line.width = Pt(1.5)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    return box


def add_title_bar(slide, title, page, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)
    add_text(slide, "車禍精神慰撫金 AI 預測系統 ｜ 資料科學期末專題",
             Inches(0.3), Inches(7.05), Inches(8), Inches(0.4), size=10, color=GREY)
    add_text(slide, f"{page} / {total}",
             Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.4),
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


TOTAL = 15

# ============== Slide 1：封面 ==============
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_text(s, "⚖️", Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5),
         size=80, align=PP_ALIGN.CENTER, color=WHITE)
add_text(s, "車禍精神慰撫金", Inches(0.5), Inches(2.7), Inches(12.3), Inches(1),
         size=48, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
add_text(s, "AI 預測系統", Inches(0.5), Inches(3.6), Inches(12.3), Inches(1),
         size=44, bold=True, align=PP_ALIGN.CENTER, color=ORANGE)
add_text(s, "用 7,050 筆真實判決，教 AI 幫你算「該賠多少 + 怎麼少賠」",
         Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
         size=20, align=PP_ALIGN.CENTER, color=WHITE)
add_text(s, "資料科學期末專題", Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
         size=16, align=PP_ALIGN.CENTER, color=LIGHT)

# ============== Slide 2：一句話介紹 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🐒 一句話介紹（猴子也懂版）", 2, TOTAL)

add_box(s, Inches(1.5), Inches(1.4), Inches(10.3), Inches(2.2), fill=LIGHT)
add_text(s, "如果你發生車禍⋯⋯",
         Inches(1.7), Inches(1.6), Inches(10), Inches(0.6),
         size=22, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "「我大概要賠多少錢？」\n「有沒有辦法少賠一點？」",
         Inches(1.7), Inches(2.2), Inches(10), Inches(1.4),
         size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, "👇 我們做了一個網站，AI 直接回答這兩個問題",
         Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
         size=20, align=PP_ALIGN.CENTER, color=GREY)

add_box(s, Inches(1.5), Inches(4.7), Inches(4.8), Inches(2.1),
        fill=RGBColor(0xE8, 0xF0, 0xFE), line=BLUE)
add_text(s, "🔮", Inches(1.5), Inches(4.8), Inches(4.8), Inches(0.7),
         size=36, align=PP_ALIGN.CENTER)
add_text(s, "功能 1：預測賠多少",
         Inches(1.5), Inches(5.55), Inches(4.8), Inches(0.5),
         size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "輸入案件條件 → 預測金額",
         Inches(1.5), Inches(6.05), Inches(4.8), Inches(0.5),
         size=14, color=GREY, align=PP_ALIGN.CENTER)

add_box(s, Inches(7.05), Inches(4.7), Inches(4.8), Inches(2.1),
        fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE)
add_text(s, "💡", Inches(7.05), Inches(4.8), Inches(4.8), Inches(0.7),
         size=36, align=PP_ALIGN.CENTER)
add_text(s, "功能 2：教你少賠一點",
         Inches(7.05), Inches(5.55), Inches(4.8), Inches(0.5),
         size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "模擬各種策略 → 算出可省金額",
         Inches(7.05), Inches(6.05), Inches(4.8), Inches(0.5),
         size=14, color=GREY, align=PP_ALIGN.CENTER)

# ============== Slide 3：為什麼做？ ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🎯 為什麼要做這個？", 3, TOTAL)

add_text(s, "❌ 問題：精神慰撫金沒有公式",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.6),
         size=24, bold=True, color=RED)
add_box(s, Inches(0.7), Inches(1.9), Inches(12), Inches(1.5), fill=LIGHT)
add_text(s,
    "法律只說「依情節輕重酌定」⋯⋯ 法官憑經驗判。\n"
    "同樣車禍，A 法官判 30 萬，B 法官可能判 80 萬，當事人完全猜不到。",
    Inches(0.9), Inches(2.05), Inches(11.7), Inches(1.3), size=18, color=NAVY)

add_text(s, "✅ 我們的解法：讓 AI 從過去判決找出規律",
         Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
         size=24, bold=True, color=GREEN)
add_box(s, Inches(0.7), Inches(4.4), Inches(12), Inches(2.2),
        fill=RGBColor(0xE8, 0xF5, 0xE9))
add_text(s,
    "1️⃣ 把過去 7,050 件車禍判決全部餵給電腦\n"
    "2️⃣ 電腦自己學：什麼條件 → 通常判多少錢\n"
    "3️⃣ 未來你輸入案情，電腦就能說：「大概賠 XX 元」",
    Inches(0.9), Inches(4.55), Inches(11.7), Inches(2),
    size=18, color=NAVY)

# ============== Slide 4：資料從哪來 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "📂 資料從哪來？", 4, TOTAL)

add_text(s, "來源：司法院裁判書開放資料",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.6),
         size=22, bold=True, color=NAVY)
add_text(s, "（所有法院判決都會公開上網，任何人都可下載）",
         Inches(0.7), Inches(1.75), Inches(12), Inches(0.5),
         size=14, color=GREY)

# 流程圖：5 個方塊
ys = Inches(2.7); h = Inches(1.5)
steps = [
    ("📦", "RAR 壓縮檔", "1996~2026 全部判決", BLUE),
    ("📄", "解壓成 JSON", "一個 JSON = 一份判決書", BLUE),
    ("🔍", "篩選車禍案件", "只留民事 + 損賠 + 車禍", ORANGE),
    ("✂️", "用正則擷取金額", "醫療費／看護費／慰撫金⋯⋯", ORANGE),
    ("📊", "7,050 筆乾淨表格", "可以餵給 AI 了", GREEN),
]
w = Inches(2.4); gap = Inches(0.15)
total_w = w * 5 + gap * 4
start = (SW - total_w) / 2
for i, (emoji, title, sub, color) in enumerate(steps):
    x = start + (w + gap) * i
    add_box(s, x, ys, w, h, fill=LIGHT, line=color)
    add_text(s, emoji, x, ys + Inches(0.1), w, Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)
    add_text(s, title, x, ys + Inches(0.7), w, Inches(0.4),
             size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, sub, x, ys + Inches(1.05), w, Inches(0.5),
             size=10, color=GREY, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow_x = x + w
        add_text(s, "▶", arrow_x - Inches(0.05), ys + Inches(0.55), gap + Inches(0.1),
                 Inches(0.5), size=14, color=color, align=PP_ALIGN.CENTER)

add_box(s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.6),
        fill=RGBColor(0xFF, 0xF8, 0xE1))
add_text(s, "📝 篩選條件（為什麼是 7,050 筆？）",
         Inches(0.9), Inches(5.1), Inches(11.7), Inches(0.4),
         size=16, bold=True, color=NAVY)
add_text(s,
    "• 民國 112~114 年（最新判決）   • 地方法院民事案件   • 案由含「損害賠償」\n"
    "• 全文含「車禍／交通事故」   • 慰撫金金額在 1 萬 ~ 1,000 萬（剔除離群值）",
    Inches(0.9), Inches(5.55), Inches(11.7), Inches(1),
    size=13, color=GREY)

# ============== Slide 5：7536 → 7050 資料清洗 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🧹 為什麼最後只用 7,050 筆？(資料清洗)", 5, TOTAL)

# 數字流程
add_box(s, Inches(0.7), Inches(1.2), Inches(3.6), Inches(1.5),
        fill=RGBColor(0xFC, 0xE4, 0xEC), line=RED)
add_text(s, "原始擷取", Inches(0.7), Inches(1.35), Inches(3.6), Inches(0.5),
         size=16, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "7,536 筆", Inches(0.7), Inches(1.8), Inches(3.6), Inches(0.7),
         size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)

add_text(s, "➡", Inches(4.4), Inches(1.6), Inches(0.7), Inches(0.7),
         size=36, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "刪 486 筆", Inches(4.1), Inches(2.25), Inches(1.4), Inches(0.4),
         size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_box(s, Inches(5.5), Inches(1.2), Inches(3.6), Inches(1.5), fill=LIGHT)
add_text(s, "清洗中⋯⋯", Inches(5.5), Inches(1.35), Inches(3.6), Inches(0.5),
         size=16, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "🧹", Inches(5.5), Inches(1.8), Inches(3.6), Inches(0.7),
         size=36, align=PP_ALIGN.CENTER)

add_text(s, "➡", Inches(9.2), Inches(1.6), Inches(0.7), Inches(0.7),
         size=36, color=GREY, align=PP_ALIGN.CENTER)

add_box(s, Inches(9.9), Inches(1.2), Inches(2.8), Inches(1.5),
        fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)
add_text(s, "訓練資料", Inches(9.9), Inches(1.35), Inches(2.8), Inches(0.5),
         size=16, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "7,050 筆", Inches(9.9), Inches(1.8), Inches(2.8), Inches(0.7),
         size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 保留率
add_text(s, "✅ 保留率 93.5%（資料品質提升，模型才會準）",
         Inches(0.7), Inches(2.85), Inches(12), Inches(0.5),
         size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 4 種被刷掉的資料
add_text(s, "🗑️ 哪些資料被刷掉了？",
         Inches(0.7), Inches(3.5), Inches(12), Inches(0.5),
         size=18, bold=True, color=NAVY)

reasons = [
    ("1️⃣", "慰撫金欄位缺失", "判決書中根本沒寫慰撫金金額 → 沒有 Y 就無法訓練 AI", RED),
    ("2️⃣", "金額過低 (< 1 萬)", "可能是正則擷取錯誤（例如把「112 年」抓成 112 元）", ORANGE),
    ("3️⃣", "金額過高 (> 1,000 萬)", "可能誤抓到訴訟費或請求總額（不是法官實際判賠的）", ORANGE),
    ("4️⃣", "其他費用 > 5,000 萬", "醫療費／看護費／工作損失明顯離譜的擷取錯誤", ORANGE),
]
y = Inches(4.05); h = Inches(0.6)
for emoji, title, desc, color in reasons:
    add_box(s, Inches(0.7), y, Inches(12), h, fill=LIGHT, line=color)
    add_text(s, emoji, Inches(0.85), y + Inches(0.1), Inches(0.5), Inches(0.4),
             size=18, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.4), y + Inches(0.13), Inches(3.5), Inches(0.4),
             size=14, bold=True, color=color)
    add_text(s, desc, Inches(5.0), y + Inches(0.15), Inches(7.5), Inches(0.4),
             size=12, color=GREY)
    y += h + Inches(0.08)

# 重點 Quote
add_box(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.45),
        fill=RGBColor(0xFF, 0xF8, 0xE1))
add_text(s, "💡 資料科學金句：「Garbage in, garbage out」── 模型品質 = 訓練資料品質",
         Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============== Slide 6：表格欄位意義 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "📋 每一筆資料代表什麼？", 6, TOTAL)

add_text(s, "每一列 = 一份車禍判決，欄位的意思：",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         size=18, color=GREY)

cols = [
    ("Injury_Level", "傷亡程度", "1=傷害 2=重傷 3=死亡", BLUE),
    ("Medical_Fee", "醫療費 (元)", "受傷花的醫藥費", BLUE),
    ("Care_Fee", "看護費 (元)", "需要人照顧的費用", BLUE),
    ("Work_Loss", "工作損失 (元)", "因傷沒上班的薪資", BLUE),
    ("Drunk", "是否酒駕", "0=否 1=是", ORANGE),
    ("Fault_Ratio", "原告過失 %", "受害人自己也有錯的比例", ORANGE),
    ("Court_XXX", "法院別", "看哪一個法院判的", GREY),
    ("Mental_Damage ⭐", "精神慰撫金 (元)", "👉 這就是我們要預測的目標！", RED),
]

# 兩欄
y0 = Inches(1.9); row_h = Inches(0.55)
for i, (en, zh, desc, color) in enumerate(cols):
    col = i % 2; row = i // 2
    x = Inches(0.7) + col * Inches(6.2)
    y = y0 + row * (row_h + Inches(0.08))
    add_box(s, x, y, Inches(6.0), row_h, fill=LIGHT)
    add_text(s, en, x + Inches(0.15), y + Inches(0.08),
             Inches(1.9), Inches(0.4), size=12, bold=True, color=color)
    add_text(s, zh, x + Inches(2.05), y + Inches(0.08),
             Inches(1.3), Inches(0.4), size=12, bold=True, color=NAVY)
    add_text(s, desc, x + Inches(3.4), y + Inches(0.1),
             Inches(2.5), Inches(0.4), size=10, color=GREY)

add_box(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6),
        fill=RGBColor(0xFF, 0xEB, 0xEE))
add_text(s,
    "💡 AI 的任務：看左邊那些欄位（X），猜出 Mental_Damage（Y）",
    Inches(0.9), Inches(6.4), Inches(11.7), Inches(0.5),
    size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ============== Slide 7：怎麼挑特徵 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🔧 資料怎麼整理？(前處理)", 7, TOTAL)

items = [
    ("1️⃣", "缺失值處理", "醫療費／看護費沒寫 → 補 0（沒請求就是 0 元）", BLUE),
    ("2️⃣", "去除離群值", "慰撫金 < 1萬 或 > 1000萬 的整筆刪除（避免怪資料汙染）", ORANGE),
    ("3️⃣", "對數轉換 Log", "金額差異大（幾萬~幾百萬），取 Log 讓數字「壓扁」一點", GREEN),
    ("4️⃣", "傷亡編碼", "「傷害／重傷／死亡」→ 1／2／3（電腦只認數字）", BLUE),
    ("5️⃣", "法院 One-Hot", "70+ 個法院 → 變成 70+ 個 0/1 欄位", ORANGE),
]

y = Inches(1.3); h = Inches(0.95)
for emoji, title, desc, color in items:
    add_box(s, Inches(0.7), y, Inches(12), h, fill=LIGHT, line=color)
    add_text(s, emoji, Inches(0.85), y + Inches(0.2), Inches(0.7), Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.6), y + Inches(0.12), Inches(3), Inches(0.5),
             size=18, bold=True, color=color)
    add_text(s, desc, Inches(1.6), y + Inches(0.5), Inches(10), Inches(0.5),
             size=14, color=GREY)
    y += h + Inches(0.1)

# ============== Slide 8：模型是什麼 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🌲 用什麼 AI 模型？Random Forest（隨機森林）", 8, TOTAL)

add_box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.5), fill=LIGHT)
add_text(s, "🐒 猴子版解釋：",
         Inches(0.9), Inches(1.3), Inches(11.7), Inches(0.5),
         size=18, bold=True, color=NAVY)
add_text(s,
    "問 100 個專家（決策樹），每個專家給出自己的答案 → 取平均當最終答案。\n"
    "為什麼比一個專家好？因為「眾人智慧」會抵銷個別專家的偏見和錯誤。",
    Inches(0.9), Inches(1.75), Inches(11.7), Inches(1),
    size=14, color=GREY)

add_text(s, "為什麼選 Random Forest？",
         Inches(0.7), Inches(2.95), Inches(12), Inches(0.5),
         size=20, bold=True, color=NAVY)

reasons = [
    ("✔", "處理表格資料超強", "我們的資料就是 Excel 一樣的表格"),
    ("✔", "不怕特徵單位不同", "金額（百萬）和過失比例（%）混在一起也沒事"),
    ("✔", "可看出特徵重要性", "可以告訴你「哪個欄位最影響判決」"),
    ("✔", "不容易過擬合", "比單一決策樹穩定，準度也比較高"),
]
y = Inches(3.6)
for icon, title, desc in reasons:
    add_box(s, Inches(0.7), y, Inches(12), Inches(0.65),
            fill=RGBColor(0xE8, 0xF5, 0xE9))
    add_text(s, icon, Inches(0.9), y + Inches(0.13), Inches(0.4), Inches(0.4),
             size=20, bold=True, color=GREEN)
    add_text(s, title, Inches(1.3), y + Inches(0.15), Inches(3.5), Inches(0.4),
             size=15, bold=True, color=NAVY)
    add_text(s, desc, Inches(4.9), y + Inches(0.17), Inches(7.7), Inches(0.4),
             size=13, color=GREY)
    y += Inches(0.75)

# ============== Slide 9：訓練流程 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🏋️ AI 怎麼學會的？(訓練流程)", 9, TOTAL)

# 流程圖
y = Inches(1.5); w = Inches(2.3); h = Inches(1.6); gap = Inches(0.25)
steps = [
    ("📊\n7,050 筆\n資料", LIGHT, NAVY),
    ("✂️\n切 8:2\n訓練／測試", RGBColor(0xE8, 0xF0, 0xFE), BLUE),
    ("🌲\n訓練\nRandom Forest", RGBColor(0xE8, 0xF5, 0xE9), GREEN),
    ("📝\n用測試集\n打分數", RGBColor(0xFF, 0xF3, 0xE0), ORANGE),
    ("💾\n存模型\nrf_model.pkl", RGBColor(0xFC, 0xE4, 0xEC), RED),
]
total_w = w * 5 + gap * 4
start = (SW - total_w) / 2
for i, (txt, fill, color) in enumerate(steps):
    x = start + (w + gap) * i
    add_box(s, x, y, w, h, fill=fill, line=color)
    add_text(s, txt, x, y + Inches(0.2), w, Inches(1.3),
             size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "➡", x + w - Inches(0.05), y + Inches(0.6), gap + Inches(0.1),
                 Inches(0.5), size=20, color=GREY, align=PP_ALIGN.CENTER)

# 訓練細節
add_box(s, Inches(0.7), Inches(3.5), Inches(12), Inches(3.3), fill=LIGHT)
add_text(s, "🔍 訓練細節（給老師看）",
         Inches(0.9), Inches(3.6), Inches(11.7), Inches(0.5),
         size=18, bold=True, color=NAVY)
add_text(s,
    "• 切資料：80% (5,640 筆) 訓練  +  20% (1,410 筆) 測試\n"
    "• 隨機種子 random_state=42（保證每次跑結果一樣）\n"
    "• 超參數：100 棵樹、最大深度 10、葉節點最少 5 筆（防止過擬合）\n"
    "• 預測目標：log(慰撫金)，預測完再用 exp 轉回原始金額\n"
    "• 評估指標：MAE（平均誤差幾元）+ R²（解釋變異量）\n"
    "• 比較基準：Null Model（不管什麼案件都猜訓練集平均值）",
    Inches(0.9), Inches(4.15), Inches(11.7), Inches(2.6),
    size=14, color=NAVY)

# ============== Slide 10：模型成績 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🏆 AI 學得怎麼樣？(模型成績)", 10, TOTAL)

# 比較卡片
add_box(s, Inches(0.7), Inches(1.3), Inches(5.8), Inches(2.5),
        fill=RGBColor(0xFC, 0xE4, 0xEC), line=RED)
add_text(s, "❌ Null Model（盲猜）",
         Inches(0.9), Inches(1.45), Inches(5.4), Inches(0.5),
         size=18, bold=True, color=RED)
add_text(s, "不管什麼案件，一律猜平均值",
         Inches(0.9), Inches(1.95), Inches(5.4), Inches(0.4),
         size=12, color=GREY)
add_text(s, "MAE = 494,720 元",
         Inches(0.9), Inches(2.5), Inches(5.4), Inches(0.7),
         size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "（平均每筆判錯約 49 萬）",
         Inches(0.9), Inches(3.2), Inches(5.4), Inches(0.4),
         size=12, color=GREY, align=PP_ALIGN.CENTER)

add_box(s, Inches(6.85), Inches(1.3), Inches(5.8), Inches(2.5),
        fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)
add_text(s, "✅ Random Forest（我們的 AI）",
         Inches(7.05), Inches(1.45), Inches(5.4), Inches(0.5),
         size=18, bold=True, color=GREEN)
add_text(s, "看過所有特徵後給出聰明答案",
         Inches(7.05), Inches(1.95), Inches(5.4), Inches(0.4),
         size=12, color=GREY)
add_text(s, "MAE = 304,969 元",
         Inches(7.05), Inches(2.5), Inches(5.4), Inches(0.7),
         size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "（平均每筆判錯約 30 萬）",
         Inches(7.05), Inches(3.2), Inches(5.4), Inches(0.4),
         size=12, color=GREY, align=PP_ALIGN.CENTER)

# 進步幅度
add_box(s, Inches(0.7), Inches(4.0), Inches(12), Inches(1.3),
        fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE)
add_text(s, "🎉 進步 38.4%！R² = 0.4618",
         Inches(0.7), Inches(4.1), Inches(12), Inches(0.8),
         size=32, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "AI 比「盲猜」聰明很多，能解釋約 46% 的判決變異",
         Inches(0.7), Inches(4.85), Inches(12), Inches(0.4),
         size=14, color=GREY, align=PP_ALIGN.CENTER)

# 為什麼不更高
add_box(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.3), fill=LIGHT)
add_text(s, "🤔 為什麼 MAE 還有 30 萬？",
         Inches(0.9), Inches(5.6), Inches(11.7), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_text(s,
    "因為精神慰撫金本來就有「法官個人心證」成分，無法 100% 預測。\n"
    "但比起「完全不知道大概多少」，30 萬的誤差已能提供合理參考區間。",
    Inches(0.9), Inches(6.0), Inches(11.7), Inches(0.8),
    size=13, color=GREY)

# ============== Slide 11：哪些因素最重要 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🔑 哪些因素最影響判決？", 11, TOTAL)

add_text(s, "AI 自己找出來的 Top 3 關鍵因素：",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         size=18, color=GREY)

# 三大條
bars = [
    ("🥇", "看護費 Care_Fee", 39, GREEN),
    ("🥈", "傷亡程度 Injury_Level", 29, BLUE),
    ("🥉", "醫療費 Medical_Fee", 19, ORANGE),
]
y = Inches(1.95); h = Inches(1.0); max_w = Inches(9)
for emoji, name, pct, color in bars:
    add_text(s, emoji, Inches(0.7), y + Inches(0.25), Inches(0.8), Inches(0.5),
             size=32, align=PP_ALIGN.CENTER)
    add_text(s, name, Inches(1.6), y + Inches(0.05), Inches(4), Inches(0.4),
             size=16, bold=True, color=NAVY)
    # bar bg
    add_box(s, Inches(1.6), y + Inches(0.5), max_w, Inches(0.45), fill=LIGHT)
    # bar fill
    bar_w = int(max_w * pct / 50)
    add_box(s, Inches(1.6), y + Inches(0.5), bar_w, Inches(0.45), fill=color)
    add_text(s, f"{pct}%", Inches(10.8), y + Inches(0.5), Inches(1.8), Inches(0.5),
             size=22, bold=True, color=color)
    y += h + Inches(0.15)

add_box(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.4),
        fill=RGBColor(0xFF, 0xF8, 0xE1))
add_text(s, "💡 解讀：法官看到「醫療費高、看護費高、傷得重」",
         Inches(0.9), Inches(5.6), Inches(11.7), Inches(0.5),
         size=16, bold=True, color=NAVY)
add_text(s, "→ 推定被害人「痛苦很深」→ 判更多慰撫金。\n"
            "  這跟我們直覺也吻合，代表 AI 學到的規律是合理的！",
         Inches(0.9), Inches(6.1), Inches(11.7), Inches(0.8),
         size=14, color=GREY)

# ============== Slide 12：功能 1 預測 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🔮 功能 1：AI 賠償預測 怎麼用？", 12, TOTAL)

add_text(s, "操作流程（三步驟）",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         size=20, bold=True, color=NAVY)

steps = [
    ("1️⃣", "填案件條件", "傷亡程度、醫療費、看護費、酒駕、過失比例...", BLUE),
    ("2️⃣", "點「進行 AI 預測」", "後端把資料送進 Random Forest 模型", ORANGE),
    ("3️⃣", "看結果", "顯示預估金額 + 參考區間 (±MAE)", GREEN),
]
y = Inches(2.0); h = Inches(1.0)
for emoji, title, desc, color in steps:
    add_box(s, Inches(0.7), y, Inches(12), h, fill=LIGHT, line=color)
    add_text(s, emoji, Inches(0.9), y + Inches(0.2), Inches(0.7), Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.7), y + Inches(0.12), Inches(4), Inches(0.5),
             size=18, bold=True, color=color)
    add_text(s, desc, Inches(1.7), y + Inches(0.55), Inches(10), Inches(0.5),
             size=13, color=GREY)
    y += h + Inches(0.15)

# 範例
add_box(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.4),
        fill=RGBColor(0xE8, 0xF0, 0xFE), line=BLUE)
add_text(s, "📋 範例：重傷 + 醫療費 10 萬 + 看護費 20 萬 + 無酒駕",
         Inches(0.9), Inches(5.6), Inches(11.7), Inches(0.5),
         size=15, bold=True, color=NAVY)
add_text(s, "👉 AI 預測：NT$ 487,984（參考區間 NT$ 183,015 ~ NT$ 792,953）",
         Inches(0.9), Inches(6.15), Inches(11.7), Inches(0.6),
         size=18, bold=True, color=BLUE)

# ============== Slide 13：功能 2 減少賠償 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "💡 功能 2：教你少賠一點 (核心亮點)", 13, TOTAL)

add_text(s, "做法：What-if 模擬",
         Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         size=20, bold=True, color=ORANGE)
add_text(s, "「固定其他條件，只改一個變數，看 AI 預測金額怎麼變」",
         Inches(0.7), Inches(1.7), Inches(12), Inches(0.5),
         size=14, color=GREY)

add_box(s, Inches(0.7), Inches(2.4), Inches(12), Inches(2.5), fill=LIGHT)
add_text(s, "🎬 舉例：原案件預測賠 80 萬",
         Inches(0.9), Inches(2.5), Inches(11.7), Inches(0.5),
         size=16, bold=True, color=NAVY)
add_text(s,
    "🅰  若主張對方有 20% 過失成功     → AI 重新預測：72 萬   省 8 萬\n"
    "🅱  若質疑看護費砍掉 50% 成功     → AI 重新預測：61 萬   省 19 萬   ⭐ 最有效\n"
    "🅲  若質疑醫療費砍掉 25% 成功     → AI 重新預測：75 萬   省 5 萬\n"
    "🅳  若質疑工作損失砍 50% 成功     → AI 重新預測：78 萬   省 2 萬",
    Inches(0.9), Inches(3.0), Inches(11.7), Inches(1.9),
    size=14, color=NAVY)

add_box(s, Inches(0.7), Inches(5.1), Inches(12), Inches(1.7),
        fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE)
add_text(s, "✨ 為什麼這很厲害？",
         Inches(0.9), Inches(5.2), Inches(11.7), Inches(0.5),
         size=16, bold=True, color=ORANGE)
add_text(s,
    "傳統律師只能說「過失比例越高賠得越少」，但說不出具體幾元。\n"
    "我們的 AI 可以「實際算出每個策略能省多少」，讓使用者排優先順序！\n"
    "👉 把預測模型「逆向使用」變成決策輔助工具",
    Inches(0.9), Inches(5.65), Inches(11.7), Inches(1.1),
    size=13, color=NAVY)

# ============== Slide 14：系統架構 ==============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "🏗️ 系統怎麼跑起來？(技術架構)", 14, TOTAL)

# 上：使用者
add_box(s, Inches(5.0), Inches(1.2), Inches(3.3), Inches(0.8),
        fill=RGBColor(0xE8, 0xF0, 0xFE), line=BLUE)
add_text(s, "👤 使用者瀏覽器",
         Inches(5.0), Inches(1.35), Inches(3.3), Inches(0.5),
         size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "⬇", Inches(6.4), Inches(2.05), Inches(0.5), Inches(0.4),
         size=20, color=GREY)

# Streamlit
add_box(s, Inches(3.5), Inches(2.5), Inches(6.3), Inches(1.0),
        fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)
add_text(s, "🎨 Streamlit 網站 (app.py + pages/)",
         Inches(3.5), Inches(2.65), Inches(6.3), Inches(0.4),
         size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "三個頁面：預測 / 減少建議 / 資料探索",
         Inches(3.5), Inches(3.05), Inches(6.3), Inches(0.4),
         size=11, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "⬇", Inches(6.4), Inches(3.55), Inches(0.5), Inches(0.4),
         size=20, color=GREY)

# predictor.py
add_box(s, Inches(3.5), Inches(4.0), Inches(6.3), Inches(1.0),
        fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE)
add_text(s, "🧠 predictor.py (模型載入與預測)",
         Inches(3.5), Inches(4.15), Inches(6.3), Inches(0.4),
         size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "load_model() / predict_amount()",
         Inches(3.5), Inches(4.55), Inches(6.3), Inches(0.4),
         size=11, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "⬇", Inches(6.4), Inches(5.05), Inches(0.5), Inches(0.4),
         size=20, color=GREY)

# 模型檔
add_box(s, Inches(1.5), Inches(5.5), Inches(4.5), Inches(1.0),
        fill=LIGHT, line=NAVY)
add_text(s, "📦 models/rf_model.pkl",
         Inches(1.5), Inches(5.65), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "已訓練的 Random Forest",
         Inches(1.5), Inches(6.05), Inches(4.5), Inches(0.4),
         size=11, color=GREY, align=PP_ALIGN.CENTER)

# 資料檔
add_box(s, Inches(7.3), Inches(5.5), Inches(4.5), Inches(1.0),
        fill=LIGHT, line=NAVY)
add_text(s, "📊 data/dataset_cleaned.csv",
         Inches(7.3), Inches(5.65), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "7,050 筆訓練資料（給 EDA 用）",
         Inches(7.3), Inches(6.05), Inches(4.5), Inches(0.4),
         size=11, color=GREY, align=PP_ALIGN.CENTER)

add_text(s, "🚀 一行指令啟動：streamlit run app.py",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============== Slide 15：總結 ==============
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
add_text(s, "✨ 總結", Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 三個亮點卡片
cards = [
    ("🎯", "解決真實問題", "精神慰撫金法官靠心證，\n民眾完全無從預估", BLUE),
    ("🤖", "用 AI 找規律", "7,050 筆判決訓練\nRandom Forest 模型", GREEN),
    ("💡", "創新亮點", "What-if 模擬\n反向給出減賠策略", ORANGE),
]
y = Inches(1.7); w = Inches(3.9); h = Inches(2.5); gap = Inches(0.2)
total_w = w * 3 + gap * 2
start = (SW - total_w) / 2
for i, (emoji, title, desc, color) in enumerate(cards):
    x = start + (w + gap) * i
    add_box(s, x, y, w, h, fill=WHITE, line=color)
    add_text(s, emoji, x, y + Inches(0.2), w, Inches(0.8),
             size=44, align=PP_ALIGN.CENTER)
    add_text(s, title, x, y + Inches(1.1), w, Inches(0.5),
             size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, y + Inches(1.6), w, Inches(0.9),
             size=13, color=GREY, align=PP_ALIGN.CENTER)

# 成果數字
add_box(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
        fill=RGBColor(0x2C, 0x4E, 0x80))
add_text(s, "📈 模型成績：MAE 304,969 元 ｜ 比盲猜進步 38.4% ｜ R² = 0.46",
         Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.6),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "🌐 互動式網站已上線可試用：streamlit run app.py",
         Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.6),
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# 致謝
add_text(s, "Thank you 🙏",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
         size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ============== 存檔 ==============
prs.save(OUTPUT)
print(f"[OK] 已產生簡報：{OUTPUT}（共 {len(prs.slides)} 頁）")
